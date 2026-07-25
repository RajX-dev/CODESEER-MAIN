from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
# Set slide dimensions to 16:9
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_dark_slide(prs, title_text, content_text):
    blank_slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(11, 15, 25) # Midnight Blue
    background.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(2))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title_text
    title_p.font.name = 'Arial'
    title_p.font.size = Pt(60)
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.font.bold = True
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(4.5))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    lines = content_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        p.text = line
        p.font.name = 'Arial'
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(220, 220, 220)

slides_data = [
    {
        "title": "N3MO: The Code Intelligence Layer",
        "content": "Deterministic code context for the modern SDLC.\n\nZero-hallucination code intelligence for enterprise,\nbanks, and governments."
    },
    {
        "title": "The Problem: AI is Blind",
        "content": "- Copilot writes bad code. RAG hallucinates.\n- PRs take too long because humans can't map\n  10,000 files in their heads.\n- Enterprise code requires certainty, not probability."
    },
    {
        "title": "The Solution: Deterministic Intelligence",
        "content": "We built a deterministic code engine.\n\n- 0 LLMs. Pure structural AST analysis.\n- 0 API calls. 100% air-gapped security.\n- 0 code stored externally."
    },
    {
        "title": "The Proof: Enterprise Scale",
        "content": "Tested on TensorFlow (36,000 files):\n\n- Indexed in <15 minutes locally.\n- 480,000+ call edges mapped.\n- Sub-second query times."
    },
    {
        "title": "Traction: Bottom-Up Adoption",
        "content": "- 8,000+ PyPI Downloads.\n- Developers are adopting N3MO to fix their\n  PR workflows today.\n- Purely organic. Zero marketing spend."
    },
    {
        "title": "The Wedge: Pull Request Reviews",
        "content": "Open-source PRs -> Team Adoption -> Enterprise\n\nWe infiltrate teams by halving the time it takes\nto review a complex PR."
    },
    {
        "title": "The Future: AI's Nervous System",
        "content": "N3MO isn't just a tool; it's the mandatory data pipeline\nfor autonomous coding agents.\n\nWhen million-agent swarms write software, they\nwill use N3MO to understand the codebase."
    },
    {
        "title": "The Ask: $500K Pre-Seed",
        "content": "- Solo technical founder (Raj Shekhar).\n- Raising $500,000 to hire elite engineers.\n- Goal: Secure first 5 enterprise pilots.\n- Mindset: Lean, focused, ramen-profitable."
    }
]

for slide in slides_data:
    add_dark_slide(prs, slide["title"], slide["content"])

output_path = "N3MO_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Successfully generated {output_path}")
