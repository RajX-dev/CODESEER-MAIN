from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Colors (Modern SaaS Dark Mode)
BG_COLOR = RGBColor(15, 23, 42)      # Slate 900
CARD_COLOR = RGBColor(30, 41, 59)    # Slate 800
ACCENT_CYAN = RGBColor(6, 182, 212)  # Cyan 500
ACCENT_PURPLE = RGBColor(168, 85, 247) # Purple 500
TEXT_MAIN = RGBColor(248, 250, 252)  # Slate 50
TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400

def add_premium_background(slide):
    # Main background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()
    
    # Top accent line (gradient simulated by two blocks)
    accent1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8), Inches(0.15))
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = ACCENT_CYAN
    accent1.line.fill.background()
    
    accent2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), 0, Inches(8), Inches(0.15))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = ACCENT_PURPLE
    accent2.line.fill.background()

def create_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_premium_background(slide)
    
    # Title
    t_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(2))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Segoe UI'
    p.font.size = Pt(88)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    # Subtitle
    s_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(1))
    sf = s_box.text_frame
    p2 = sf.paragraphs[0]
    p2.text = subtitle
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(36)
    p2.font.color.rgb = ACCENT_CYAN

def create_content_slide(title, bullets, cards=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_premium_background(slide)
    
    # Title
    t_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1.5))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Segoe UI'
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    if cards:
        # Draw cards instead of standard text
        card_width = Inches(4.3)
        card_height = Inches(5.0)
        for i, card_text in enumerate(cards):
            left = Inches(1) + (i * Inches(4.7))
            
            # Card background
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.8), card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_COLOR
            # Minimalistic border
            card.line.color.rgb = ACCENT_CYAN
            card.line.width = Pt(1.5)
            
            # Card Text
            c_box = slide.shapes.add_textbox(left + Inches(0.3), Inches(3.1), card_width - Inches(0.6), card_height - Inches(0.6))
            c_tf = c_box.text_frame
            c_tf.word_wrap = True
            
            lines = card_text.split('\n')
            for j, line in enumerate(lines):
                if j == 0:
                    cp = c_tf.paragraphs[0]
                    cp.font.color.rgb = ACCENT_CYAN
                    cp.font.size = Pt(32)
                    cp.font.bold = True
                else:
                    cp = c_tf.add_paragraph()
                    cp.font.color.rgb = TEXT_MAIN
                    cp.font.size = Pt(24)
                cp.text = line
                cp.font.name = 'Segoe UI'
    else:
        # Standard bullets with custom styling
        c_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(5))
        c_tf = c_box.text_frame
        c_tf.word_wrap = True
        
        for i, b in enumerate(bullets):
            if i == 0:
                p = c_tf.paragraphs[0]
            else:
                p = c_tf.add_paragraph()
            p.text = "■  " + b  
            p.font.name = 'Segoe UI'
            p.font.size = Pt(38)
            p.font.color.rgb = TEXT_MUTED
            p.space_after = Pt(36)

# Slide 1
create_title_slide("N3MO", "The Deterministic Code Intelligence Layer.")

# Slide 2
create_content_slide(
    "The Problem: AI is Blind",
    [],
    cards=[
        "Hallucinations\n\nCopilot writes bad code because it relies on probabilistic RAG instead of structure.",
        "Slow PRs\n\nHuman reviewers cannot mentally map the blast radius across a 10,000-file monorepo.",
        "Security Risks\n\nBanks and defense sectors cannot exfiltrate code to external APIs. They need certainty."
    ]
)

# Slide 3
create_content_slide(
    "The Solution: Deterministic Intelligence",
    [
        "0 LLMs. Pure structural AST graph analysis.",
        "0 API calls. 100% air-gapped security out of the box.",
        "0 code stored externally. Ultimate privacy for enterprises."
    ]
)

# Slide 4
create_content_slide(
    "The Proof: Enterprise Scale",
    [],
    cards=[
        "36,000 Files\n\nThe massive TensorFlow monorepo indexed completely locally.",
        "14 Minutes\n\nFrom zero to full graph extraction in under 15 minutes.",
        "< 200 ms\n\nSub-second query times across 480,000+ recursive call edges."
    ]
)

# Slide 5
create_content_slide(
    "Traction: Bottom-Up Adoption",
    [
        "8,000+ PyPI Downloads.",
        "Adopted organically by developers to fix their daily PR workflows.",
        "Zero marketing spend. Pure product-led growth."
    ]
)

# Slide 6
create_content_slide(
    "The Wedge: Pull Requests",
    [],
    cards=[
        "1. Open Source\n\nOSS Maintainers use N3MO to map dependencies instantly.",
        "2. Team Adoption\n\nMaintainers bring it to their corporate engineering teams.",
        "3. Enterprise CI\n\nSecurity teams mandate it for air-gapped CI/CD pipelines."
    ]
)

# Slide 7
create_content_slide(
    "Competitive Moat vs All-in-Ones",
    [],
    cards=[
        "Speed & Focus\n\nUnlike heavy platforms (Repowise, Sonar), N3MO does one thing perfectly: deterministic blast-radius mapping.",
        "Zero Hallucination\n\nCompetitors use LLMs for analysis. N3MO uses pure math and ASTs—no drift, no embedding costs.",
        "Frictionless CI/CD\n\nN3MO targets the exact pain point: automatic PR webhooks with zero config. No local servers required."
    ]
)

# Slide 8
create_content_slide(
    "The Ask: $500K Pre-Seed",
    [
        "Solo technical founder: Raj Shekhar (Architect of N3MO).",
        "Goal: Secure first 5 enterprise pilots & hire elite founding engineers.",
        "Mindset: Ultra-lean, obsessed with engineering, ramen-profitable."
    ]
)

prs.save("N3MO_Pitch_Deck_Premium.pptx")
print("Saved N3MO_Pitch_Deck_Premium.pptx")
